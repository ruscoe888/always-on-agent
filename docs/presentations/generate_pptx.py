"""
Generate three .pptx files from the always-on ops agent presentations.
Matches the Batman dark theme as closely as PowerPoint allows.
Run: python3 docs/generate_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Colours ──────────────────────────────────────────────────────────────────
BG        = RGBColor(0x08, 0x08, 0x08)
CARD      = RGBColor(0x10, 0x10, 0x10)
YELLOW    = RGBColor(0xFF, 0xE3, 0x00)
GREEN     = RGBColor(0x00, 0xE6, 0x76)
ORANGE    = RGBColor(0xFF, 0x6B, 0x35)
TRIAGE    = RGBColor(0xFF, 0x6B, 0x35)
CORR      = RGBColor(0x29, 0xB6, 0xF6)
COMP      = RGBColor(0xBA, 0x68, 0xC8)
P0        = RGBColor(0xFF, 0x17, 0x44)
P1        = RGBColor(0xFF, 0x6B, 0x35)
P2        = RGBColor(0xFF, 0xE3, 0x00)
P3        = RGBColor(0x29, 0xB6, 0xF6)
WHITE     = RGBColor(0xD8, 0xD8, 0xD8)
MUTED     = RGBColor(0x55, 0x55, 0x55)
MUTED2    = RGBColor(0x77, 0x77, 0x77)

SLIDE_W   = Inches(13.33)
SLIDE_H   = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_layout(prs):
    return prs.slide_layouts[6]  # completely blank

def bg(slide):
    """Fill slide background solid dark."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

def box(slide, x, y, w, h, fill_color=None, border_color=None, border_pt=0.75):
    """Add a rectangle shape."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.width = Pt(border_pt) if border_color else Pt(0)
    if border_color:
        shape.line.color.rgb = border_color
    else:
        shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape

def accent_bar(slide, x, y, h, color):
    """Thin vertical accent bar."""
    bar = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.06), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

def txt(slide, text, x, y, w, h,
        size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        font_name="Calibri", wrap=True):
    """Add a text box."""
    tf_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tf_box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font_name
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return tf_box

def heading(slide, text, x, y, w, color=YELLOW, size=40):
    return txt(slide, text, x, y, w, 1.2, size=size, bold=True, color=color, font_name="Calibri")

def eyebrow(slide, text, x, y, w):
    return txt(slide, text.upper(), x, y, w, 0.35, size=10, color=MUTED, font_name="Courier New")

def body(slide, text, x, y, w, h, size=14, color=MUTED2):
    return txt(slide, text, x, y, w, h, size=size, color=color, wrap=True)

def bullet_block(slide, items, x, y, w, accent_color, size=13):
    """Render a list of (icon, text) bullet items."""
    for i, (icon, text) in enumerate(items):
        yi = y + i * 0.52
        txt(slide, icon, x, yi, 0.35, 0.5, size=16, color=accent_color)
        txt(slide, text, x + 0.38, yi, w - 0.38, 0.5, size=size, color=MUTED2)

def divider(slide, x, y, w):
    """Horizontal rule."""
    line = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.012))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x22)
    line.line.fill.background()

def progress_bar(slide, current, total, accent):
    """Yellow/green/orange progress indicator at top."""
    w = 13.33 * (current / total)
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(w), Inches(0.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

def footer(slide, deck_name, slide_num, total, accent):
    """Bottom chrome bar."""
    box(slide, 0, 7.3, 13.33, 0.2, fill_color=RGBColor(0x0a, 0x0a, 0x0a))
    txt(slide, deck_name, 0.3, 7.32, 8, 0.18, size=8, color=MUTED, font_name="Courier New")
    txt(slide, f"{slide_num} / {total}", 12.5, 7.32, 0.8, 0.18, size=8, color=MUTED,
        font_name="Courier New", align=PP_ALIGN.RIGHT)

def section_card(slide, x, y, w, h, title, body_text, accent):
    box(slide, x, y, w, h, fill_color=CARD, border_color=RGBColor(0x25,0x25,0x25))
    accent_bar(slide, x, y, h, accent)
    txt(slide, title, x+0.18, y+0.1, w-0.25, 0.35, size=13, bold=True, color=accent)
    txt(slide, body_text, x+0.18, y+0.48, w-0.25, h-0.55, size=11, color=MUTED2)


# ═══════════════════════════════════════════════════════════════════════════════
# DECK 1 — CLIENT (Batman yellow)
# ═══════════════════════════════════════════════════════════════════════════════

def build_client(path):
    prs = new_prs()
    TOTAL = 7
    ACC   = YELLOW
    DNAME = "Wayne Enterprises — Ops Intelligence  |  CONFIDENTIAL"

    # ── Slide 1: Cover ─────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout(prs))
    bg(sl); progress_bar(sl, 1, TOTAL, ACC); footer(sl, DNAME, 1, TOTAL, ACC)
    eyebrow(sl, "Confidential — Client Briefing", 0.7, 0.6, 9)
    heading(sl, "Your ops team is fighting fires\ninstead of preventing them.", 0.7, 1.0, 9, color=ACC, size=44)
    body(sl, ("Every hour your team spends manually triaging incidents, chasing deploy correlations, "
              "and auditing vendor contracts is an hour not spent on engineering that matters."),
         0.7, 3.4, 7.5, 1.2, size=16)
    divider(sl, 0.7, 4.85, 5)
    txt(sl, "Presented by  Wayne Enterprises — Applied AI", 0.7, 5.05, 8, 0.4, size=12, color=MUTED)

    # ── Slide 2: Situation ─────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout(prs))
    bg(sl); progress_bar(sl, 2, TOTAL, ACC); footer(sl, DNAME, 2, TOTAL, ACC)
    eyebrow(sl, "The Situation", 0.7, 0.6, 9)
    heading(sl, "Sound familiar?", 0.7, 0.95, 9, color=ACC, size=36)
    items = [
        ("🔔", "Incidents land in a queue. Your team reads each one, searches for a runbook, and manually decides severity. 20–40 min per incident."),
        ("🔗", "Deploy correlations are guesswork. 'Did a recent release cause this?' requires a human to cross-reference timestamps — always late."),
        ("📄", "Compliance drift goes unnoticed. Contracts are signed, filed, never re-audited. Violations accumulate silently until an incident forces a review."),
        ("🌙", "Nights and weekends are unmonitored. Issues surface at 2am, sit untriaged until standup. By then the customer has called the CEO."),
    ]
    for i, (icon, text) in enumerate(items):
        yi = 1.9 + i * 1.1
        box(sl, 0.7, yi, 11.9, 1.0, fill_color=RGBColor(0x12,0x07,0x07), border_color=RGBColor(0x40,0x10,0x10))
        txt(sl, icon, 0.95, yi+0.22, 0.5, 0.6, size=18)
        txt(sl, text, 1.55, yi+0.12, 10.8, 0.8, size=13, color=MUTED2)

    # ── Slide 3: The Answer ────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout(prs))
    bg(sl); progress_bar(sl, 3, TOTAL, ACC); footer(sl, DNAME, 3, TOTAL, ACC)
    eyebrow(sl, "The Answer — Lead With The Conclusion", 0.7, 0.6, 9)
    heading(sl, "An agent that never sleeps, never misses, never forgets.", 0.7, 0.95, 11, color=ACC, size=34)
    box(sl, 0.7, 2.3, 11.9, 1.2, fill_color=RGBColor(0x10,0x0E,0x00), border_color=RGBColor(0x50,0x46,0x00))
    accent_bar(sl, 0.7, 2.3, 1.2, ACC)
    txt(sl, ("We built an always-on ops agent that automatically triages every incident, correlates it with "
             "recent deployments, audits your vendor contracts, and files actionable GitHub Issues — every four hours, "
             "with zero human intervention required."),
        0.95, 2.42, 11.4, 1.0, size=14, color=WHITE)
    stats = [("18", "Issues filed\nautomatically", ACC),
             ("0",  "Permission\nprompts", P0),
             ("4h", "Max detection\ntime", CORR),
             ("9",  "Compliance\nviolations found", COMP)]
    for i, (val, lbl, col) in enumerate(stats):
        xi = 0.7 + i * 3.1
        box(sl, xi, 3.75, 2.9, 1.4, fill_color=CARD, border_color=RGBColor(0x25,0x25,0x25))
        txt(sl, val, xi+0.2, 3.82, 2.5, 0.75, size=44, bold=True, color=col)
        txt(sl, lbl, xi+0.2, 4.55, 2.5, 0.55, size=11, color=MUTED)
    body(sl, ("It runs whether you are in a meeting, asleep, or on holiday. New data in the repo triggers "
              "a full re-analysis. Findings already filed are skipped — no noise, no duplicates."),
         0.7, 5.35, 11.9, 0.8, size=13)

    # ── Slide 4: How It Works ──────────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout(prs))
    bg(sl); progress_bar(sl, 4, TOTAL, ACC); footer(sl, DNAME, 4, TOTAL, ACC)
    eyebrow(sl, "Supporting Argument 1 — Architecture", 0.7, 0.6, 9)
    heading(sl, "Two layers. One agent.", 0.7, 0.95, 9, color=ACC, size=34)
    flow = ["⏰ Cron Fires\nEvery 4 hours.\nZero infra.",
            "📥 Git Pull\nFetches latest\nissues & data.",
            "🧠 CLAUDE.md\nAll rules &\nlogic. In git.",
            "🔍 Dedup\nSkips already-\nreported issues.",
            "📋 Issues Filed\nNet-new findings\nonly. Labelled."]
    for i, step in enumerate(flow):
        xi = 0.5 + i * 2.55
        box(sl, xi, 2.1, 2.35, 1.5, fill_color=CARD, border_color=RGBColor(0x25,0x25,0x25))
        txt(sl, step, xi+0.12, 2.18, 2.15, 1.35, size=11, color=MUTED2)
        if i < 4:
            txt(sl, "→", xi+2.35, 2.7, 0.25, 0.4, size=14, color=MUTED)
    cards = [("1", "Ops Triage", "Assigns P0–P3 severity, matches runbooks, recommends on-call team and specific fix steps.", TRIAGE),
             ("2", "Incident Correlation", "Cross-references incidents with deploys within 48 hours. Identifies causal chains. Recommends rollbacks.", CORR),
             ("3", "Compliance Audit", "Checks every vendor contract against your 7-rule policy. Flags violations by severity with exact clause references.", COMP)]
    for i, (num, title, desc, col) in enumerate(cards):
        xi = 0.5 + i * 4.3
        section_card(sl, xi, 3.85, 4.1, 2.9, f"{num}  {title}", desc, col)

    # ── Slide 5: Proof ─────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout(prs))
    bg(sl); progress_bar(sl, 5, TOTAL, ACC); footer(sl, DNAME, 5, TOTAL, ACC)
    eyebrow(sl, "Supporting Argument 2 — Evidence From a Real Run", 0.7, 0.6, 9)
    heading(sl, "What it found in the first run.", 0.7, 0.95, 9, color=ACC, size=34)
    incidents = [("P0", "PROD-4521: PaymentService NPE — Correlated to payment-service v4.8.2, deployed 14 min prior. Rollback recommended.", P0),
                 ("P1", "PROD-4487: Acme Corp Checkout — Causal chain: feature flag + NPE bug. 2,400 users affected.", P1),
                 ("P2", "PROD-4519: Slow Uploads — Signing-service URL TTL reduced 3600s→300s. Rollback available.", P2),
                 ("P1", "PROD-4533: Reporting DB Timeouts — 3 enterprise customers, 18% error rate, no deploy correlation.", P1)]
    for i, (sev, desc, col) in enumerate(incidents):
        yi = 2.0 + i * 0.72
        box(sl, 0.7, yi, 0.55, 0.58, fill_color=RGBColor(0x12,0x08,0x08), border_color=col)
        txt(sl, sev, 0.78, yi+0.1, 0.45, 0.4, size=11, bold=True, color=col, font_name="Courier New")
        txt(sl, desc, 1.38, yi+0.08, 5.2, 0.5, size=11, color=MUTED2)
    txt(sl, "COMPLIANCE VIOLATIONS", 7.5, 1.85, 5.5, 0.3, size=9, color=MUTED, font_name="Courier New")
    violations = [("CRITICAL", "Sirius Storage: EU data in Malaysia only — no residency guarantee", P0),
                  ("CRITICAL", "Sirius Storage: Breach notification 'in due course' — no 72h window", P0),
                  ("HIGH",     "Sirius Storage: 5-year lock-in, no termination for convenience", P1),
                  ("CRITICAL", "Acme: Breach notification 96h (policy requires 72h)", P0),
                  ("MEDIUM",   "+ 5 further violations across both contracts", MUTED2)]
    for i, (sev, desc, col) in enumerate(violations):
        yi = 2.2 + i * 0.88
        box(sl, 7.5, yi, 5.5, 0.75, fill_color=RGBColor(0x0c,0x08,0x10), border_color=RGBColor(0x25,0x18,0x30))
        accent_bar(sl, 7.5, yi, 0.75, col)
        txt(sl, sev, 7.72, yi+0.04, 1.5, 0.28, size=9, color=col, font_name="Courier New")
        txt(sl, desc, 7.72, yi+0.32, 5.1, 0.38, size=11, color=MUTED2)

    # ── Slide 6: Dashboard ─────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout(prs))
    bg(sl); progress_bar(sl, 6, TOTAL, ACC); footer(sl, DNAME, 6, TOTAL, ACC)
    eyebrow(sl, "Supporting Argument 3 — Visibility", 0.7, 0.6, 9)
    heading(sl, "Command-centre visibility for your entire team.", 0.7, 0.95, 9, color=ACC, size=32)
    feats = [("◆", "Real-time stats — open issues, P0s, new today, compliance violations at a glance"),
             ("◆", "Five drilldown charts — by type, by day, by week, by month; click any segment for the underlying issues"),
             ("◆", "Zero infrastructure — a single HTML file reading live from the GitHub API; share the URL, done"),
             ("◆", "Automatic updates — every agent run adds new data; the dashboard reflects it on next refresh")]
    for i, (icon, text) in enumerate(feats):
        yi = 2.1 + i * 0.75
        txt(sl, icon, 0.7, yi, 0.4, 0.6, size=14, color=ACC)
        txt(sl, text, 1.15, yi+0.04, 5.8, 0.65, size=13, color=MUTED2)
    # mini dashboard preview
    box(sl, 7.5, 1.85, 5.5, 5.2, fill_color=RGBColor(0x0a,0x0a,0x0a), border_color=RGBColor(0x25,0x25,0x25))
    txt(sl, "WAYNE ENTERPRISES — COMMAND CENTER", 7.65, 1.95, 5.2, 0.3, size=7, color=MUTED, font_name="Courier New")
    stat_data = [("20", "OPEN", ACC, 7.55, 2.38), ("1", "P0", P0, 10.1, 2.38),
                 ("20", "TODAY", CORR, 7.55, 3.45), ("9", "COMPLIANCE", COMP, 10.1, 3.45)]
    for val, lbl, col, xi, yi in stat_data:
        box(sl, xi, yi, 2.35, 0.85, fill_color=RGBColor(0x11,0x11,0x11), border_color=RGBColor(0x22,0x22,0x22))
        accent_bar(sl, xi, yi, 0.85, col)
        txt(sl, val, xi+0.18, yi+0.02, 1.8, 0.48, size=30, bold=True, color=col)
        txt(sl, lbl, xi+0.18, yi+0.54, 1.8, 0.26, size=8, color=MUTED, font_name="Courier New")
    bars = [("TRIAGE", 0.38, TRIAGE), ("CORRELATION", 0.28, CORR), ("COMPLIANCE", 0.64, COMP)]
    for i, (lbl, pct, col) in enumerate(bars):
        yi = 4.5 + i * 0.5
        txt(sl, lbl, 7.65, yi, 1.5, 0.38, size=8, color=col, font_name="Courier New")
        bw = pct * 3.8
        box(sl, 9.3, yi+0.06, 3.8, 0.22, fill_color=RGBColor(0x15,0x15,0x15))
        box(sl, 9.3, yi+0.06, bw,  0.22, fill_color=col)

    # ── Slide 7: Next Steps ────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout(prs))
    bg(sl); progress_bar(sl, 7, TOTAL, ACC); footer(sl, DNAME, 7, TOTAL, ACC)
    eyebrow(sl, "Recommendation", 0.7, 0.6, 9)
    heading(sl, "What we do next, together.", 0.7, 0.95, 9, color=ACC, size=34)
    steps = [("Step 1", "Fork this pattern into your own repo — point the agent at your real incident data"),
             ("Step 2", "Extend CLAUDE.md with your runbooks, your severity definitions, your compliance policy"),
             ("Step 3", "Add Slack or PagerDuty webhook so P0s page on-call immediately"),
             ("Step 4", "Schedule the 7-day cron renewal — or promote to a GitHub Action for full autonomy")]
    for i, (badge, text) in enumerate(steps):
        yi = 2.1 + i * 0.78
        box(sl, 0.7, yi, 1.1, 0.55, fill_color=RGBColor(0x18,0x16,0x00), border_color=RGBColor(0x50,0x46,0x00))
        txt(sl, badge, 0.78, yi+0.1, 1.0, 0.38, size=11, bold=True, color=ACC)
        txt(sl, text, 1.95, yi+0.1, 5.5, 0.45, size=13, color=MUTED2)
    deliverables = ["A working always-on agent in your GitHub repo — running right now",
                    "20 auto-filed GitHub Issues: triage, correlation, and compliance",
                    "A live ops dashboard requiring no infrastructure or backend",
                    "A versioned, extensible CLAUDE.md that you own",
                    "Zero ongoing infrastructure cost — runs inside Claude Code"]
    txt(sl, "WHAT YOU WALK AWAY WITH", 7.5, 1.95, 5.5, 0.3, size=9, color=MUTED, font_name="Courier New")
    for i, d in enumerate(deliverables):
        yi = 2.35 + i * 0.72
        box(sl, 7.5, yi, 5.5, 0.6, fill_color=CARD, border_color=RGBColor(0x1e,0x1e,0x1e))
        txt(sl, "✅", 7.65, yi+0.1, 0.4, 0.4, size=13)
        txt(sl, d, 8.12, yi+0.1, 4.75, 0.45, size=11, color=MUTED2)
    box(sl, 7.5, 6.0, 5.5, 0.85, fill_color=RGBColor(0x10,0x0E,0x00), border_color=RGBColor(0x50,0x46,0x00))
    txt(sl, ("The agent paid for itself before this presentation ended. It found 9 compliance violations "
             "your legal team hasn't seen yet."), 7.65, 6.08, 5.2, 0.72, size=11, color=WHITE)

    prs.save(path)
    print(f"Saved: {path}")


# ═══════════════════════════════════════════════════════════════════════════════
# DECK 2 — TECHNICAL (green)
# ═══════════════════════════════════════════════════════════════════════════════

def build_technical(path):
    prs  = new_prs()
    TOTAL = 8
    ACC   = GREEN
    DNAME = "Always-On Ops Agent — Engineering Deep Dive  |  INTERNAL"

    def sl_base(n):
        sl = prs.slides.add_slide(blank_layout(prs))
        bg(sl); progress_bar(sl, n, TOTAL, ACC); footer(sl, DNAME, n, TOTAL, ACC)
        return sl

    # Slide 1: Overview
    sl = sl_base(1)
    eyebrow(sl, "Internal Engineering Brief — Always-On Ops Agent", 0.7, 0.6, 12)
    heading(sl, "What we built, how it works,\nwhat to know before you extend it.", 0.7, 0.95, 11, color=ACC, size=36)
    body(sl, ("This is not a demo. This is the actual system, running every 4 hours. "
              "This deck covers architecture decisions, constraints, tradeoffs, and known rough edges "
              "so you can extend it without breaking it."), 0.7, 3.0, 10, 0.9, size=15)
    meta = [("REPO", "ruscoe888/always-on-agent"),
            ("RUNTIME", "Claude Code durable cron (CronCreate)"),
            ("SCHEDULE", "Every 4h at :17 — off-peak, avoids fleet contention"),
            ("OUTPUT", "GitHub Issues via gh CLI"),
            ("PERSISTENCE", ".claude/scheduled_tasks.json — 7-day TTL, must renew")]
    for i, (k, v) in enumerate(meta):
        txt(sl, k, 0.7, 4.1 + i*0.45, 2.2, 0.4, size=10, color=MUTED, font_name="Courier New")
        txt(sl, v, 3.0, 4.1 + i*0.45, 9.5, 0.4, size=12, color=MUTED2)

    # Slide 2: Architecture
    sl = sl_base(2)
    eyebrow(sl, "Architecture", 0.7, 0.6, 9)
    heading(sl, "Two-layer design: cron prompt + CLAUDE.md", 0.7, 0.95, 11, color=ACC, size=30)
    body(sl, ("Key decision: keep the cron prompt minimal (~150 words), push all domain logic into CLAUDE.md. "
              "Claude Code reads CLAUDE.md automatically when it enters the repo directory — so the agent "
              "inherits the full protocol without the cron prompt growing unbounded."), 0.7, 2.05, 11.5, 0.8, size=13)
    arch = [("CronCreate (durable)\n~150-word prompt. Fires every 4h. Tells Claude to cd, git pull, follow CLAUDE.md.", ACC),
            ("git -C /repo pull\nFetches latest data. Repo = single source of truth.", GREEN),
            ("CLAUDE.md loaded\nSeverity rubric, runbook keywords, correlation logic, compliance rules, dedup, output format.", GREEN),
            ("gh issue list (dedup)\nQueries ops-agent issues. Matches PROD-IDs. Skips duplicates.", GREEN),
            ("gh issue create\nOne issue per net-new finding. Labels + structured body with source refs.", GREEN)]
    for i, (text, col) in enumerate(arch):
        xi = 0.5 + i * 2.55
        box(sl, xi, 3.1, 2.35, 1.65, fill_color=CARD, border_color=RGBColor(0x18,0x2e,0x18))
        accent_bar(sl, xi, 3.1, 1.65, col)
        txt(sl, text, xi+0.18, 3.18, 2.1, 1.5, size=10, color=MUTED2)
        if i < 4:
            txt(sl, "→", xi+2.35, 3.78, 0.22, 0.4, size=12, color=MUTED)
    box(sl, 0.7, 5.0, 11.9, 0.8, fill_color=RGBColor(0x06,0x10,0x06), border_color=RGBColor(0x18,0x30,0x18))
    txt(sl, ("Why not a monolithic prompt? Cron prompts are string literals — unversioned, unreviewed. "
             "CLAUDE.md is versioned in git, diffable, PR-reviewable. Changing agent behaviour = a commit."),
        0.95, 5.1, 11.4, 0.65, size=12, color=WHITE)

    # Slide 3: Permissions
    sl = sl_base(3)
    eyebrow(sl, "Permissions Model", 0.7, 0.6, 9)
    heading(sl, "Why the agent ran without prompting.", 0.7, 0.95, 9, color=ACC, size=30)
    body(sl, ("Claude Code prompts for any Bash command not in an explicit allowlist. In an unattended cron run, "
              "a prompt = a hang. We solved this at two layers: global settings and project settings."),
         0.7, 2.0, 8, 0.7, size=13)
    layers = [("GLOBAL  ~/.claude/settings.json",
               "Allows git pull, gh issue *, gh label *, gh api, cat, ls, find, echo.\nApplies to all Claude Code sessions on this machine, including cron runs.", GREEN),
              ("PROJECT  .claude/settings.json (committed)",
               "Mirrors the same allowlist. Documents requirements for a different machine.\nCommitted to the repo — visible and reviewable.", GREEN),
              ("GAP  — Novel Bash patterns will still prompt",
               "If you add new shell commands to CLAUDE.md, add them to the allowlist first.\nThe Read tool doesn't need allowlisting — it's not Bash.", P0)]
    for i, (title, desc, col) in enumerate(layers):
        yi = 2.95 + i * 1.1
        box(sl, 0.7, yi, 7.0, 0.95, fill_color=CARD, border_color=RGBColor(0x20,0x20,0x20))
        accent_bar(sl, 0.7, yi, 0.95, col)
        txt(sl, title, 0.95, yi+0.05, 6.6, 0.32, size=10, bold=True, color=col, font_name="Courier New")
        txt(sl, desc, 0.95, yi+0.4, 6.6, 0.52, size=11, color=MUTED2)
    box(sl, 0.7, 6.25, 7.0, 0.6, fill_color=RGBColor(0x10,0x05,0x05), border_color=RGBColor(0x40,0x10,0x10))
    txt(sl, ("Security: Bash(gh issue create*) allows any gh issue create invocation. "
             "If the agent reads untrusted data that could influence tool calls, this is a prompt injection surface. "
             "Current data (controlled JSON/markdown repo) is low-risk."), 0.9, 6.32, 6.6, 0.5, size=11, color=MUTED2)
    # code snippet
    box(sl, 8.0, 2.0, 5.0, 4.9, fill_color=RGBColor(0x03,0x03,0x03), border_color=RGBColor(0x18,0x30,0x18))
    code = ('~/.claude/settings.json\n\n'
            '{\n'
            '  "permissions": {\n'
            '    "allow": [\n'
            '      "Bash(git pull*)",\n'
            '      "Bash(git -C * pull*)",\n'
            '      "Bash(gh issue list*)",\n'
            '      "Bash(gh issue create*)",\n'
            '      "Bash(gh label list*)",\n'
            '      "Bash(gh label create*)",\n'
            '      "Bash(gh api*)",\n'
            '      "Bash(cat *)",\n'
            '      "Bash(ls *)",\n'
            '      "Bash(find *)",\n'
            '      "Bash(echo *)"\n'
            '    ]\n'
            '  }\n'
            '}')
    txt(sl, code, 8.15, 2.1, 4.7, 4.7, size=10, color=MUTED2, font_name="Courier New")

    # Slide 4: CLAUDE.md
    sl = sl_base(4)
    eyebrow(sl, "The Agent Contract", 0.7, 0.6, 9)
    heading(sl, "CLAUDE.md is the source of truth.", 0.7, 0.95, 9, color=ACC, size=30)
    body(sl, ("CLAUDE.md is loaded automatically by Claude Code before any run. "
              "Think of it as a system prompt that lives in version control. "
              "Extending the agent = committing a change to this file."), 0.7, 2.0, 8, 0.7, size=13)
    sections = [("Section 1", "Ops Triage Protocol", "Severity rubric (P0–P3), runbook keyword matching, on-call routing. Same input → same severity label.", ACC),
                ("Section 2", "Incident Correlation", "48-hour window, service-area matching, causal chain detection. Depends on ISO 8601 timestamps.", CORR),
                ("Section 3", "Compliance Audit", "7-rule policy check per contract. References compliance-policy.md at runtime — policy changes reflect immediately.", COMP),
                ("Section 4", "Dedup + Output", "Issue title format, label schema, dedup query logic, run summary. Deterministic — manually testable.", MUTED2)]
    for i, (tag, title, desc, col) in enumerate(sections):
        yi = 2.95 + i * 0.95
        box(sl, 0.7, yi, 7.0, 0.82, fill_color=CARD, border_color=RGBColor(0x20,0x20,0x20))
        accent_bar(sl, 0.7, yi, 0.82, col)
        txt(sl, tag, 0.95, yi+0.04, 1.3, 0.26, size=9, color=MUTED, font_name="Courier New")
        txt(sl, title, 0.95, yi+0.3, 3.5, 0.35, size=13, bold=True, color=col)
        txt(sl, desc, 4.5, yi+0.12, 3.05, 0.62, size=11, color=MUTED2)
    box(sl, 0.7, 6.75, 7.0, 0.5, fill_color=RGBColor(0x06,0x10,0x06), border_color=RGBColor(0x18,0x30,0x18))
    txt(sl, ("Extending: New runbook → update Section 1 keywords. New compliance rule → append to compliance-policy.md + Section 3. "
             "New output channel → add Section 4 + allowlist the Bash commands."), 0.9, 6.82, 6.6, 0.38, size=11, color=WHITE)
    box(sl, 8.0, 2.0, 5.0, 5.3, fill_color=RGBColor(0x03,0x03,0x03), border_color=RGBColor(0x18,0x30,0x18))
    code2 = ('CLAUDE.md — severity rubric excerpt\n\n'
             '## Task 1: Ops Triage\n\n'
             'Severity:\n'
             '- P0: Multiple customers, revenue\n'
             '       impact, active outage\n'
             '- P1: Single enterprise customer\n'
             '       or service degraded\n'
             '- P2: Intermittent, limited impact\n'
             '- P3: Feature requests, non-urgent\n\n'
             'Runbook keywords:\n'
             '- auth-502-windows.md\n'
             '  → 502, login, auth, pool\n'
             '- cdn-upload-latency.md\n'
             '  → upload, CDN, slow, image\n'
             '- payment-service-degraded.md\n'
             '  → PaymentService, NPE, checkout')
    txt(sl, code2, 8.15, 2.1, 4.7, 5.1, size=10, color=MUTED2, font_name="Courier New")

    # Slide 5: Data Schema
    sl = sl_base(5)
    eyebrow(sl, "Data Layer", 0.7, 0.6, 9)
    heading(sl, "What the agent reads — and what it expects.", 0.7, 0.95, 11, color=ACC, size=30)
    schemas = [
        ("issues/PROD-XXXX.json",
         '{\n  "id": "PROD-4521",\n  "title": "...",\n  "status": "open",\n  "severity": null,\n'
         '  "opened_at": "ISO8601",  ← critical\n  "reporter": "email",\n  "body": "...",\n  "comments": []\n}',
         "opened_at must be ISO 8601.\nMissing = silent correlation failure."),
        ("deploys/recent.json",
         '{\n  "deploys": [{\n    "service": "payment-service",\n    "version": "v4.8.2",\n'
         '    "deployed_at": "ISO8601",\n    "commit": "9f4a1c8",\n    "rollback_available": true,\n'
         '    "last_known_good": "v4.8.1"\n  }]\n}',
         "Manually maintained today.\nNext: CI/CD appends on every deploy."),
        ("compliance-policy.md",
         '## Data residency\nEU data must stay in EU.\n\n## Breach notification\n≤72 hours from detection.\n\n'
         '## Termination\n≤90 days notice for convenience.\n\n## Liability cap\n≥12 months. Breach not excluded.\n...',
         "Read at runtime each run.\nUpdate the file = instant policy change.")
    ]
    for i, (title, code_text, note) in enumerate(schemas):
        xi = 0.5 + i * 4.3
        txt(sl, title, xi, 2.05, 4.1, 0.3, size=9, color=MUTED, font_name="Courier New")
        box(sl, xi, 2.38, 4.1, 3.8, fill_color=RGBColor(0x03,0x03,0x03), border_color=RGBColor(0x18,0x2e,0x18))
        txt(sl, code_text, xi+0.12, 2.48, 3.85, 3.6, size=9.5, color=MUTED2, font_name="Courier New")
        txt(sl, note, xi, 6.28, 4.1, 0.65, size=11, color=MUTED)

    # Slide 6: Tradeoffs
    sl = sl_base(6)
    eyebrow(sl, "Design Decisions — Honest Assessment", 0.7, 0.6, 9)
    heading(sl, "What we chose and what we gave up.", 0.7, 0.95, 11, color=ACC, size=30)
    rows = [("Runtime", "Claude Code durable cron", "GitHub Actions", "Zero infra, zero YAML", "7-day TTL. Requires CC running."),
            ("Output", "GitHub Issues via gh", "Commit markdown reports", "Issues are actionable + closeable", "Issues accumulate. No auto-close."),
            ("Logic location", "CLAUDE.md in repo", "Embedded in cron prompt", "Versioned, reviewable, auto-loaded", "Must cd to correct directory."),
            ("Deduplication", "gh issue list title-match", "Local state / database", "GitHub is the source of truth", "Slow at >200 issues. Substring match."),
            ("Deploy data", "Manual JSON file", "Live CI/CD API query", "No auth tokens needed in agent", "Stale if deploys not appended manually.")]
    hdrs = ["DECISION", "CHOSEN", "ALTERNATIVE", "WHY WE CHOSE THIS", "WHAT WE GAVE UP"]
    cols_w = [1.9, 2.1, 2.0, 2.8, 2.8]
    x_starts = [0.5, 2.42, 4.54, 6.56, 9.38]
    for j, (hdr, xj) in enumerate(zip(hdrs, x_starts)):
        txt(sl, hdr, xj, 2.05, cols_w[j], 0.3, size=8, color=MUTED, font_name="Courier New")
    divider(sl, 0.5, 2.38, 12.3)
    for i, row in enumerate(rows):
        yi = 2.52 + i * 0.88
        if i % 2 == 0:
            box(sl, 0.5, yi, 12.3, 0.84, fill_color=RGBColor(0x0c,0x0c,0x0c))
        for j, (cell, xj) in enumerate(zip(row, x_starts)):
            col = MUTED2 if j < 3 else (P0 if j == 4 else MUTED2)
            txt(sl, cell, xj, yi+0.1, cols_w[j]-0.1, 0.72, size=11, color=col)

    # Slide 7: Extension Points
    sl = sl_base(7)
    eyebrow(sl, "Known Gaps and Extension Points", 0.7, 0.6, 9)
    heading(sl, "What to build next.", 0.7, 0.95, 9, color=ACC, size=30)
    gaps = [("1", "7-day cron TTL", "Auto-expires — requires manual renewal. Fix: self-renewal cron or GitHub Action."),
            ("2", "Deploy data is manual", "deploys/recent.json is hand-maintained. CI/CD should append on every deploy."),
            ("3", "Dedup degrades at scale", "gh issue list --limit 200. Beyond 200: add pagination or a state file."),
            ("4", "No P0 escalation path", "Agent files a GitHub Issue for a P0. Needs a Slack/PagerDuty webhook.")]
    txt(sl, "GAPS", 0.7, 2.0, 5.8, 0.3, size=9, color=P0, font_name="Courier New")
    for i, (num, title, desc) in enumerate(gaps):
        yi = 2.35 + i * 1.0
        box(sl, 0.7, yi, 5.8, 0.88, fill_color=CARD, border_color=RGBColor(0x20,0x20,0x20))
        txt(sl, num, 0.88, yi+0.1, 0.35, 0.6, size=22, bold=True, color=RGBColor(0x28,0x28,0x28))
        txt(sl, title, 1.28, yi+0.06, 5.0, 0.32, size=12, bold=True, color=WHITE)
        txt(sl, desc, 1.28, yi+0.44, 5.0, 0.38, size=11, color=MUTED2)
    exts = [("A", "New runbooks", "Add markdown to runbooks/ + keywords to CLAUDE.md Task 1. Done."),
            ("B", "New compliance rules", "Append to compliance-policy.md + add check to CLAUDE.md Task 3."),
            ("C", "New output channels", "Add Task 4 to CLAUDE.md with channel logic. Allowlist the Bash commands."),
            ("D", "Real-time trigger", "Wire GitHub webhook on issue creation to trigger the agent immediately via gh workflow dispatch.")]
    txt(sl, "EXTENSION POINTS", 7.1, 2.0, 5.8, 0.3, size=9, color=ACC, font_name="Courier New")
    for i, (ltr, title, desc) in enumerate(exts):
        yi = 2.35 + i * 1.0
        box(sl, 7.1, yi, 5.8, 0.88, fill_color=CARD, border_color=RGBColor(0x18,0x2e,0x18))
        accent_bar(sl, 7.1, yi, 0.88, ACC)
        txt(sl, ltr, 7.3, yi+0.1, 0.35, 0.6, size=22, bold=True, color=RGBColor(0x18,0x30,0x18))
        txt(sl, title, 7.65, yi+0.06, 5.1, 0.32, size=12, bold=True, color=ACC)
        txt(sl, desc, 7.65, yi+0.44, 5.1, 0.38, size=11, color=MUTED2)

    # Slide 8: Open Questions
    sl = sl_base(8)
    eyebrow(sl, "For Discussion", 0.7, 0.6, 9)
    heading(sl, "Open questions for the team.", 0.7, 0.95, 9, color=ACC, size=30)
    questions = [
        ("1", "Who owns CLAUDE.md?",
         "Any engineer can PR a change. Unreviewed changes to the severity rubric or dedup logic could cause mis-triage or duplicates. Should it require a designated reviewer?"),
        ("2", "Where should the cron live long-term?",
         "Currently tied to one machine. For production reliability: GitHub Actions, Lambda, or managed Claude API scheduled invocation. Each has different auth and cost profiles."),
        ("3", "How do we validate agent output quality?",
         "Severity assignments are LLM outputs — plausible but not guaranteed correct. Should there be a human review gate before P0 issues page on-call?"),
        ("4", "Is the prompt injection surface acceptable?",
         "Agent reads issue bodies from a controlled repo. If anyone can write to the repo, they can craft an issue body that attempts to manipulate tool calls."),
        ("5", "What's the rollback story?",
         "If the agent files 50 incorrect issues due to a bad CLAUDE.md change, there's no bulk-close command. Should we add a cleanup script?"),
    ]
    for i, (num, q, detail) in enumerate(questions):
        yi = 2.05 + i * 1.02
        box(sl, 0.7, yi, 12.3, 0.9, fill_color=CARD, border_color=RGBColor(0x20,0x20,0x20))
        txt(sl, num, 0.88, yi+0.1, 0.4, 0.7, size=22, bold=True, color=RGBColor(0x18,0x30,0x18))
        txt(sl, q, 1.38, yi+0.06, 4.2, 0.35, size=13, bold=True, color=WHITE)
        txt(sl, detail, 5.7, yi+0.06, 7.1, 0.78, size=11, color=MUTED2)

    prs.save(path)
    print(f"Saved: {path}")


# ═══════════════════════════════════════════════════════════════════════════════
# DECK 3 — OPS TEAM (orange)
# ═══════════════════════════════════════════════════════════════════════════════

def build_ops(path):
    prs  = new_prs()
    TOTAL = 7
    ACC   = ORANGE
    DNAME = "Always-On Ops Agent — Ops Team Briefing  |  INTERNAL"

    def sl_base(n):
        sl = prs.slides.add_slide(blank_layout(prs))
        bg(sl); progress_bar(sl, n, TOTAL, ACC); footer(sl, DNAME, n, TOTAL, ACC)
        return sl

    # Slide 1: What Changed
    sl = sl_base(1)
    eyebrow(sl, "Ops Team Briefing — Always-On Agent", 0.7, 0.6, 9)
    heading(sl, "Your triage queue now has a first responder.", 0.7, 0.95, 11, color=ACC, size=36)
    body(sl, ("An automated agent now runs every four hours and does the first pass on every open incident — "
              "severity assignment, runbook lookup, deploy correlation, and compliance checks. "
              "This brief explains what it does, what it creates, and what you need to action."),
         0.7, 2.3, 10, 0.9, size=15)
    labels = [("ops-agent", RGBColor(0x1d,0x76,0xdb)), ("agent-triage", TRIAGE),
              ("agent-correlation", CORR), ("agent-compliance", COMP)]
    txt(sl, "LABELS IT CREATES", 0.7, 3.55, 6, 0.28, size=8, color=MUTED, font_name="Courier New")
    for i, (lbl, col) in enumerate(labels):
        xi = 0.7 + i * 3.1
        box(sl, xi, 3.88, 2.85, 0.42, fill_color=RGBColor(0x0c,0x0c,0x0c), border_color=col)
        txt(sl, lbl, xi+0.12, 3.96, 2.65, 0.3, size=11, color=col, font_name="Courier New")
    schedule = [("WHERE", "github.com/ruscoe888/always-on-agent → Issues → filter: label:ops-agent"),
                ("RUNS AT", "00:17 · 04:17 · 08:17 · 12:17 · 16:17 · 20:17 UTC (every 4 hours)")]
    for i, (k, v) in enumerate(schedule):
        txt(sl, k, 0.7, 4.65+i*0.5, 1.5, 0.38, size=9, color=MUTED, font_name="Courier New")
        txt(sl, v, 2.3, 4.65+i*0.5, 10.0, 0.38, size=13, color=MUTED2)

    # Slide 2: Before / After
    sl = sl_base(2)
    eyebrow(sl, "What Changed For Your Team", 0.7, 0.6, 9)
    heading(sl, "Your morning looked different yesterday.", 0.7, 0.95, 11, color=ACC, size=32)
    box(sl, 0.5, 2.1, 6.1, 5.15, fill_color=RGBColor(0x10,0x06,0x06), border_color=RGBColor(0x35,0x10,0x10))
    txt(sl, "BEFORE — Manual First Response", 0.7, 2.18, 5.7, 0.3, size=9, color=P0, font_name="Courier New")
    before = [("🕐", "Alert fires. Engineer reads it, decides if real. 5–15 min before anyone acts."),
              ("📖", "Search Confluence for the runbook. Hope it's up to date."),
              ("🔍", "Manually check deploy logs: 'did anything go out?'. Another 10 minutes."),
              ("📋", "Vendor compliance? Reviewed once at signing, never touched again."),
              ("🌙", "2am incident sits untriaged until standup. Customer has escalated.")]
    for i, (icon, text) in enumerate(before):
        txt(sl, icon, 0.7, 2.55+i*0.82, 0.4, 0.7, size=14)
        txt(sl, text, 1.15, 2.6+i*0.82, 5.2, 0.7, size=12, color=MUTED2)
    box(sl, 6.75, 2.1, 6.1, 5.15, fill_color=RGBColor(0x06,0x10,0x06), border_color=RGBColor(0x10,0x30,0x10))
    txt(sl, "AFTER — Agent First Response", 6.95, 2.18, 5.7, 0.3, size=9, color=GREEN, font_name="Courier New")
    after = [("⚡", "Within 4 hours of any new incident, the agent has read it, assigned severity, and filed a GitHub Issue."),
             ("📋", "The issue body already contains the matching runbook and specific fix steps."),
             ("🔗", "If a deploy in the last 48 hours could have caused it, the agent has already identified it."),
             ("⚖️", "Vendor contracts are re-audited on every run. New violations filed automatically."),
             ("✅", "Your team arrives to a pre-triaged queue. You decide what to act on, not what to read.")]
    for i, (icon, text) in enumerate(after):
        txt(sl, icon, 6.95, 2.55+i*0.82, 0.4, 0.7, size=14)
        txt(sl, text, 7.4, 2.6+i*0.82, 5.2, 0.7, size=12, color=MUTED2)

    # Slide 3: Reading an Issue
    sl = sl_base(3)
    eyebrow(sl, "How To Use This In Your Daily Workflow", 0.7, 0.6, 9)
    heading(sl, "Reading an agent-filed issue.", 0.7, 0.95, 9, color=ACC, size=30)
    fields = [("TITLE", "[Agent Triage] in the title = agent filed it. Severity (P1) is in the title — scan the list without opening each issue.", ACC),
              ("SEVERITY", "Agent's assigned severity with reasoning. Treat as an informed first opinion — override it if wrong and leave a comment.", CORR),
              ("RUNBOOK", "Matched runbook file + fix steps. If 'No matching runbook found': incident is novel, diagnose manually, consider adding a runbook.", GREEN),
              ("ESCALATION", "Which on-call team to page. Check this first for P0/P1 — agent has identified the right team from the service area.", MUTED2)]
    for i, (key, desc, col) in enumerate(fields):
        yi = 2.1 + i * 0.88
        box(sl, 0.7, yi, 6.5, 0.78, fill_color=CARD, border_color=RGBColor(0x20,0x20,0x20))
        accent_bar(sl, 0.7, yi, 0.78, col)
        txt(sl, key, 0.95, yi+0.06, 1.4, 0.3, size=9, color=col, font_name="Courier New")
        txt(sl, desc, 0.95, yi+0.38, 6.0, 0.38, size=11, color=MUTED2)
    box(sl, 0.7, 5.7, 6.5, 0.6, fill_color=RGBColor(0x12,0x08,0x02), border_color=RGBColor(0x40,0x22,0x08))
    txt(sl, ("Important: The agent files issues — it does not close them. "
             "Once resolved, close the GitHub Issue manually. This keeps the dashboard accurate."),
        0.9, 5.78, 6.1, 0.48, size=11, color=WHITE)
    # Issue anatomy preview
    box(sl, 7.5, 2.0, 5.5, 4.5, fill_color=RGBColor(0x0a,0x0a,0x0a), border_color=RGBColor(0x30,0x1a,0x05))
    txt(sl, "[Agent Triage] PROD-4521: NullPointerException\nin PaymentService at checkout (P0)", 7.65, 2.08, 5.2, 0.65, size=11, bold=True, color=ORANGE)
    anatomy = [("Severity:", "P0 — Active outage. p99 8.2s. Error rate +340%.\n2,400 checkouts affected. Revenue impact ongoing."),
               ("Runbook:", "✅ runbooks/payment-service-degraded.md"),
               ("Action:", "1. payment-service v4.8.2 deployed 14 min before\n2. NPE at PaymentService.java:142 — null in guest checkout\n3. Rollback to v4.8.1 ← RECOMMENDED"),
               ("Escalate:", "#payments-oncall — page immediately")]
    yi = 2.85
    for key, val in anatomy:
        txt(sl, key, 7.65, yi, 1.2, 0.38, size=9, color=MUTED, font_name="Courier New")
        txt(sl, val, 8.9, yi, 3.95, 0.8, size=10, color=MUTED2)
        yi += 0.82

    # Slide 4: Severity Guide
    sl = sl_base(4)
    eyebrow(sl, "Reference — Severity Definitions", 0.7, 0.6, 9)
    heading(sl, "What P0–P3 means in practice.", 0.7, 0.95, 11, color=ACC, size=30)
    sev_rows = [
        ("P0", P0, "Active revenue-impacting outage. Multiple customers or entire service down.",
         "PROD-4521: PaymentService NPE. 2,400 checkouts broken.", "Page on-call immediately. Incident bridge.", "Ack 15min. Update every 30min."),
        ("P1", P1, "Single enterprise customer affected, or significant degradation with workaround.",
         "PROD-4487: Acme Corp checkout broken. CTO on call.", "Page on-call + notify CSM within 1h.", "Ack 1h. Fix/workaround 4h."),
        ("P2", P2, "Intermittent degradation. Limited blast radius. Self-recovering.",
         "PROD-4498: Login 502s in 30s windows, self-resolving.", "Assign to next sprint / on-call team.", "Ack 24h. Fix in 5 biz days."),
        ("P3", P3, "Feature request, minor inconvenience, no customer impact.",
         "PROD-4506: Parallel batch job support — internal.", "Route to product backlog. No urgency.", "No SLA."),
    ]
    hdrs2 = ["SEV", "MEANS", "EXAMPLE", "YOUR RESPONSE", "SLA"]
    cw = [0.6, 2.4, 3.0, 2.9, 1.8]
    xs = [0.5, 1.18, 3.65, 6.72, 9.65]
    for j, (hdr, xj) in enumerate(zip(hdrs2, xs)):
        txt(sl, hdr, xj, 2.08, cw[j], 0.28, size=8, color=MUTED, font_name="Courier New")
    divider(sl, 0.5, 2.4, 11.0)
    for i, (sev, col, means, example, resp, sla) in enumerate(sev_rows):
        yi = 2.55 + i * 1.02
        if i % 2 == 0:
            box(sl, 0.5, yi, 11.0, 0.95, fill_color=RGBColor(0x0c,0x0c,0x0c))
        txt(sl, sev, xs[0], yi+0.3, cw[0], 0.42, size=14, bold=True, color=col, font_name="Courier New")
        txt(sl, means,   xs[1], yi+0.08, cw[1], 0.82, size=11, color=MUTED2)
        txt(sl, example, xs[2], yi+0.08, cw[2], 0.82, size=11, color=MUTED2)
        txt(sl, resp,    xs[3], yi+0.08, cw[3], 0.82, size=11, color=MUTED2)
        txt(sl, sla,     xs[4], yi+0.08, cw[4], 0.82, size=11, color=MUTED2)
    box(sl, 0.5, 6.7, 11.0, 0.55, fill_color=RGBColor(0x12,0x08,0x02), border_color=RGBColor(0x40,0x22,0x08))
    txt(sl, ("If you disagree with severity: edit the issue title and leave a comment. "
             "The agent deduplicates by PROD-ID so your override is preserved on re-runs."),
        0.7, 6.78, 10.6, 0.4, size=11, color=WHITE)

    # Slide 5: What It Found
    sl = sl_base(5)
    eyebrow(sl, "First Run Results — What The Agent Found", 0.7, 0.6, 9)
    heading(sl, "20 issues filed. Zero missed.", 0.7, 0.95, 9, color=ACC, size=30)
    # Triage col
    txt(sl, "TRIAGE (5)", 0.5, 2.0, 4.0, 0.28, size=9, color=TRIAGE, font_name="Courier New")
    triage_items = [("P0", "PaymentService NPE — checkouts broken us-east-1 + eu-west-2", P0),
                    ("P1", "Acme Corp checkout spinner — 2,400 users, CTO escalation", P1),
                    ("P1", "Reporting DB timeouts — 3 enterprise customers, 18% errors", P1),
                    ("P2", "Login 502 windows — intermittent, self-resolving", P2),
                    ("P2", "Slow image uploads — 30–60s on iOS and Android", P2)]
    for i, (sev, desc, col) in enumerate(triage_items):
        yi = 2.38 + i * 0.82
        box(sl, 0.5, yi, 0.55, 0.66, fill_color=RGBColor(0x12,0x08,0x08), border_color=col)
        txt(sl, sev, 0.57, yi+0.16, 0.45, 0.36, size=10, bold=True, color=col, font_name="Courier New")
        txt(sl, desc, 1.12, yi+0.12, 3.25, 0.52, size=11, color=MUTED2)
    # Correlation col
    txt(sl, "CORRELATIONS (4)", 4.6, 2.0, 4.0, 0.28, size=9, color=CORR, font_name="Courier New")
    corrs = [("PROD-4521 ← payment-service v4.8.2", "Deployed 14 min before. Rollback available."),
             ("PROD-4487 ← CAUSAL CHAIN", "payment-service v4.8.2 + tenant-config v3.2.1.\nBoth rollbacks available."),
             ("PROD-4519 ← signing-service v2.1.4", "URL TTL 3600s→300s. Rollback available."),
             ("PROD-4498 ← auth-service v6.0.0", "Redis migration. No rollback available.")]
    for i, (title, desc) in enumerate(corrs):
        yi = 2.38 + i * 1.08
        box(sl, 4.6, yi, 3.9, 0.92, fill_color=RGBColor(0x06,0x0d,0x12), border_color=RGBColor(0x10,0x25,0x35))
        txt(sl, title, 4.75, yi+0.06, 3.6, 0.32, size=10, bold=True, color=CORR)
        txt(sl, desc, 4.75, yi+0.44, 3.6, 0.44, size=11, color=MUTED2)
    # Compliance col
    txt(sl, "COMPLIANCE (9)", 8.75, 2.0, 4.2, 0.28, size=9, color=COMP, font_name="Courier New")
    viols = [("CRITICAL", "Sirius: EU data in Malaysia only", P0),
             ("CRITICAL", "Sirius: Breach 'in due course' — no 72h window", P0),
             ("HIGH",     "Sirius: 5-year lock-in, no exit clause", P1),
             ("HIGH",     "Sirius: Liability cap 3 months (need 12)", P1),
             ("CRITICAL", "Acme: Breach notification 96h (need 72h)", P0),
             ("+ 4 more", "view all in GitHub Issues", MUTED)]
    for i, (sev, desc, col) in enumerate(viols):
        yi = 2.38 + i * 0.72
        box(sl, 8.75, yi, 4.2, 0.62, fill_color=RGBColor(0x0c,0x08,0x10), border_color=RGBColor(0x20,0x15,0x28))
        accent_bar(sl, 8.75, yi, 0.62, col)
        txt(sl, sev, 8.98, yi+0.04, 1.3, 0.24, size=8, color=col, font_name="Courier New")
        txt(sl, desc, 8.98, yi+0.32, 3.8, 0.28, size=11, color=MUTED2)

    # Slide 6: Daily Workflow
    sl = sl_base(6)
    eyebrow(sl, "How To Integrate This Into Your Routine", 0.7, 0.6, 9)
    heading(sl, "A suggested daily workflow.", 0.7, 0.95, 9, color=ACC, size=30)
    steps = [("1", "Open GitHub Issues filtered by ops-agent",
              "Filter: label=ops-agent, state=open. Agent runs at 04:17 so overnight incidents are pre-triaged when you arrive."),
             ("2", "Scan for P0 and P1 first",
              "Severity is in the issue title. Any open unassigned P0 = first action before anything else."),
             ("3", "Check correlation issues for P0/P1s",
              "Filter by agent-correlation. Rollback recommendation in the body if a deploy is the cause."),
             ("4", "Check the dashboard for the broader picture",
              "dashboard/index.html — click any chart segment for the underlying issues.")]
    for i, (num, title, desc) in enumerate(steps):
        yi = 2.1 + i * 1.1
        box(sl, 0.7, yi, 6.3, 0.95, fill_color=CARD, border_color=RGBColor(0x20,0x20,0x20))
        box(sl, 0.7, yi, 0.55, 0.95, fill_color=RGBColor(0x15,0x0a,0x03), border_color=RGBColor(0x35,0x1a,0x05))
        txt(sl, num, 0.82, yi+0.2, 0.4, 0.55, size=22, bold=True, color=ACC)
        txt(sl, title, 1.38, yi+0.08, 5.4, 0.35, size=12, bold=True, color=WHITE)
        txt(sl, desc, 1.38, yi+0.52, 5.4, 0.38, size=11, color=MUTED2)
    txt(sl, "CLOSING ISSUES + LIMITS", 7.3, 2.0, 5.5, 0.28, size=9, color=MUTED, font_name="Courier New")
    closing = [("Incident resolved", "Close the GitHub Issue manually. Leave a comment with the resolution."),
               ("Compliance remediated", "Close + note contract version updated. Legal team should confirm."),
               ("Severity is wrong", "Edit the title, leave a comment. Your edit is preserved on re-runs.")]
    for i, (title, desc) in enumerate(closing):
        yi = 2.35 + i * 0.7
        box(sl, 7.3, yi, 5.5, 0.6, fill_color=CARD, border_color=RGBColor(0x20,0x20,0x20))
        txt(sl, title, 7.46, yi+0.06, 2.5, 0.28, size=11, bold=True, color=WHITE)
        txt(sl, desc, 7.46, yi+0.34, 5.2, 0.24, size=11, color=MUTED2)
    txt(sl, "WHAT THE AGENT DOES NOT DO", 7.3, 4.48, 5.5, 0.28, size=9, color=P0, font_name="Courier New")
    limits = ["Page on-call directly (webhook coming next)",
              "Close resolved issues — you do that",
              "Diagnose issues without a matching runbook",
              "Execute rollbacks — it recommends, you action",
              "Read live monitoring systems — repo data only"]
    for i, lim in enumerate(limits):
        txt(sl, "✗", 7.3, 4.82+i*0.46, 0.3, 0.38, size=11, color=P0)
        txt(sl, lim, 7.62, 4.82+i*0.46, 5.1, 0.38, size=11, color=MUTED2)

    # Slide 7: Feedback Loop
    sl = sl_base(7)
    eyebrow(sl, "Making It Better — Your Input Matters", 0.7, 0.6, 9)
    heading(sl, "The agent learns from your feedback.", 0.7, 0.95, 9, color=ACC, size=30)
    feedback = [("📖", "Wrong runbook match?", "Tell us the PROD-ID and which runbook should have matched. We update the keyword table in CLAUDE.md.", ACC),
                ("🔗", "Missed a correlation?", "Give us the PROD-ID + deploy. We check if the timestamp window or service-area matching needs updating.", CORR),
                ("⚖️", "New compliance rule?", "Legal adds to compliance-policy.md. Agent picks it up on next run. No code change needed.", COMP),
                ("📌", "New runbook?", "Add markdown to runbooks/ + keywords to CLAUDE.md Task 1. Done.", MUTED2)]
    for i, (icon, title, desc, col) in enumerate(feedback):
        yi = 2.1 + i * 1.08
        box(sl, 0.7, yi, 6.5, 0.92, fill_color=CARD, border_color=RGBColor(0x20,0x20,0x20))
        accent_bar(sl, 0.7, yi, 0.92, col)
        txt(sl, icon, 0.85, yi+0.2, 0.5, 0.55, size=18)
        txt(sl, title, 1.42, yi+0.08, 5.5, 0.32, size=13, bold=True, color=col)
        txt(sl, desc, 1.42, yi+0.5, 5.5, 0.38, size=11, color=MUTED2)
    # Quick ref
    box(sl, 7.5, 2.0, 5.5, 5.25, fill_color=RGBColor(0x0a,0x0a,0x0a), border_color=RGBColor(0x25,0x25,0x25))
    txt(sl, "QUICK REFERENCE", 7.65, 2.08, 5.2, 0.28, size=8, color=MUTED, font_name="Courier New")
    ref = [("Dashboard", "dashboard/index.html (open locally)"),
           ("Issues filter", "label:ops-agent is:open"),
           ("Agent runs", "00:17 · 04:17 · 08:17 · 12:17 · 16:17 · 20:17"),
           ("Protocol", "CLAUDE.md (root of repo)"),
           ("Policy", "compliance-policy.md (root of repo)"),
           ("Runbooks", "runbooks/ (add new ones here)"),
           ("Cron TTL", "7 days — renew before expiry"),
           ("Feedback", "Open a PR against CLAUDE.md")]
    for i, (k, v) in enumerate(ref):
        yi = 2.48 + i * 0.58
        txt(sl, k, 7.65, yi, 1.8, 0.44, size=9, color=MUTED, font_name="Courier New")
        col = ORANGE if k == "Cron TTL" else MUTED2
        txt(sl, v, 9.55, yi, 3.3, 0.44, size=11, color=col)

    prs.save(path)
    print(f"Saved: {path}")


# ── Main ──────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    out = os.path.dirname(os.path.abspath(__file__))
    build_client(    os.path.join(out, "presentation-client.pptx"))
    build_technical( os.path.join(out, "presentation-technical.pptx"))
    build_ops(       os.path.join(out, "presentation-ops.pptx"))
    print("\nAll three decks generated.")
